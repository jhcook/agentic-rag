"""Text extraction from various file formats."""
import io
import logging
import pathlib
from typing import Union

try:
    from pypdf import PdfReader  # type: ignore
except ImportError:
    PdfReader = None

try:
    from docx import Document  # type: ignore
except ImportError:
    Document = None

try:
    from bs4 import BeautifulSoup  # type: ignore
except ImportError:
    BeautifulSoup = None

try:
    import csv  # stdlib
except ImportError:
    csv = None

try:
    import openpyxl  # type: ignore
except ImportError:
    openpyxl = None

try:
    from pptx import Presentation  # type: ignore
except ImportError:
    Presentation = None

try:
    from striprtf.striprtf import rtf_to_text  # type: ignore
except ImportError:
    rtf_to_text = None

try:
    import ebooklib  # type: ignore
    from ebooklib import epub  # type: ignore
except ImportError:
    ebooklib = None
    epub = None

try:
    from PIL import Image  # type: ignore
except ImportError:
    Image = None

try:
    import pytesseract  # type: ignore
except ImportError:
    pytesseract = None

try:
    import requests  # type: ignore
except ImportError:
    requests = None

logger = logging.getLogger(__name__)


def _download_from_url(url: str) -> bytes:
    """Download content from a URL."""
    if requests is None:
        logger.warning("URL support not available. Install requests: pip install requests")
        return b""
    try:
        logger.info("Downloading from %s", url)
        response = requests.get(url, timeout=30)
        response.raise_for_status()
        return response.content
    except requests.exceptions.SSLError as ssl_exc:
        logger.error("SSL error downloading %s: %s", url, ssl_exc)
        return b"[SSL ERROR: Could not connect to %s]" % url.encode()
    except requests.exceptions.ConnectionError as conn_exc:
        logger.error("Connection error downloading %s: %s", url, conn_exc)
        return b"[CONNECTION ERROR: Could not connect to %s]" % url.encode()
    except Exception as exc:  # pylint: disable=broad-exception-caught
        logger.error("Failed to download %s: %s", url, exc)
        return b"[ERROR: Could not download %s: %s]" % (url.encode(), str(exc).encode())


def extract_text_from_bytes(content: bytes, filename: str) -> str:
    """Extract text from file content (bytes) based on filename extension."""
    suffix = pathlib.Path(filename).suffix.lower()

    if suffix == '.pdf':
        if PdfReader is None:
            logger.warning("PDF support not available. Install pypdf: pip install pypdf")
            return ""
        try:
            reader = PdfReader(io.BytesIO(content))
            text_parts = []
            for page in reader.pages:
                page_text = page.extract_text() or ""
                if page_text:
                    text_parts.append(page_text)
            joined = "\n".join(text_parts).strip()
            if not joined:
                logger.warning("No text extracted from PDF %s", filename)
                return ""
            return joined
        except Exception as exc:  # pylint: disable=broad-exception-caught
            logger.error("Failed to read PDF %s: %s", filename, exc)
            return ""

    elif suffix in ['.docx', '.doc']:
        if Document is None:
            logger.warning(
                "DOCX support not available. Install python-docx: pip install python-docx")
            return ""
        try:
            doc = Document(io.BytesIO(content))
            text_parts = [paragraph.text for paragraph in doc.paragraphs]
            return "\n".join(text_parts)
        except Exception as exc:  # pylint: disable=broad-exception-caught
            logger.error("Failed to read DOCX %s: %s", filename, exc)
            return ""

    elif suffix in ['.md', '.markdown']:
        try:
            text = content.decode('utf-8', errors='ignore')
            return text
        except Exception as exc:  # pylint: disable=broad-exception-caught
            logger.error("Failed to read Markdown %s: %s", filename, exc)
            return ""

    elif suffix == '.csv':
        try:
            decoded = content.decode('utf-8', errors='ignore')
            return decoded
        except Exception as exc:  # pylint: disable=broad-exception-caught
            logger.error("Failed to read CSV %s: %s", filename, exc)
            return ""

    elif suffix in ['.xlsx', '.xls']:
        if openpyxl is None:
            logger.warning("XLSX support not available. Install openpyxl: pip install openpyxl")
            return ""
        try:
            wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
            text_parts = []
            for sheet in wb.worksheets[:3]:  # limit to first 3 sheets to avoid bloat
                for row in sheet.iter_rows(values_only=True):
                    if row:
                        text_parts.append(" ".join("" if v is None else str(v) for v in row))
            return "\n".join(text_parts)
        except Exception as exc:  # pylint: disable=broad-exception-caught
            logger.error("Failed to read XLSX %s: %s", filename, exc)
            return ""

    elif suffix in ['.pptx', '.ppt']:
        if Presentation is None:
            logger.warning("PPTX support not available. Install python-pptx: pip install python-pptx")
            return ""
        try:
            pres = Presentation(io.BytesIO(content))
            text_parts = []
            for slide in pres.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_parts.append(shape.text)
            return "\n".join(text_parts)
        except Exception as exc:  # pylint: disable=broad-exception-caught
            logger.error("Failed to read PPTX %s: %s", filename, exc)
            return ""

    elif suffix == '.rtf':
        if rtf_to_text is None:
            logger.warning("RTF support not available. Install striprtf: pip install striprtf")
            return ""
        try:
            text = rtf_to_text(content.decode('utf-8', errors='ignore'))
            return text
        except Exception as exc:  # pylint: disable=broad-exception-caught
            logger.error("Failed to read RTF %s: %s", filename, exc)
            return ""

    elif suffix == '.epub':
        if epub is None:
            logger.warning("EPUB support not available. Install ebooklib: pip install ebooklib")
            return ""
        try:
            book = epub.read_epub(io.BytesIO(content))
            text_parts = []
            for item in book.get_items():
                if item.get_type() == ebooklib.ITEM_DOCUMENT:
                    try:
                        html_content = item.get_content().decode('utf-8', errors='ignore')
                        if BeautifulSoup:
                            soup = BeautifulSoup(html_content, 'html.parser')
                            for script in soup(["script", "style"]):
                                script.decompose()
                            text_parts.append(soup.get_text(separator='\n'))
                        else:
                            text_parts.append(html_content)
                    except Exception:  # pylint: disable=broad-exception-caught
                        continue
            return "\n".join(text_parts)
        except Exception as exc:  # pylint: disable=broad-exception-caught
            logger.error("Failed to read EPUB %s: %s", filename, exc)
            return ""

    elif suffix in ['.html', '.htm']:
        if BeautifulSoup is None:
            logger.warning(
                "HTML support not available. Install beautifulsoup4: "
                "pip install beautifulsoup4")
            return ""
        try:
            html_content = content.decode('utf-8', errors='ignore')
            soup = BeautifulSoup(html_content, 'html.parser')
            for script in soup(["script", "style"]):
                script.decompose()
            text = soup.get_text(separator='\n')
            lines = (line.strip() for line in text.splitlines())
            chunks = (line.strip() for line in lines)
            text = '\n'.join(chunk for chunk in chunks if chunk)
            return text
        except Exception as exc:  # pylint: disable=broad-exception-caught
            logger.error("Failed to read HTML %s: %s", filename, exc)
            return ""

    else:
        # Fallback: sniff for HTML even if extension is missing
        if BeautifulSoup is not None:
            try:
                prefix = content[:2048].lower()
                if b"<html" in prefix or b"<!doctype html" in prefix:
                    soup = BeautifulSoup(content, "html.parser")
                    for script in soup(["script", "style"]):
                        script.decompose()
                    text = soup.get_text(separator="\n")
                    lines = (line.strip() for line in text.splitlines())
                    chunks = (line.strip() for line in lines)
                    return "\n".join(chunk for chunk in chunks if chunk)
            except Exception as exc:  # pylint: disable=broad-exception-caught
                logger.warning("HTML sniff parse failed for %s: %s", filename, exc)

        # Fallback: simple OCR for common image types if available
        if suffix in ['.png', '.jpg', '.jpeg', '.tiff', '.bmp'] and Image is not None and pytesseract is not None:
            try:
                with Image.open(io.BytesIO(content)) as img:
                    return pytesseract.image_to_string(img)
            except Exception as exc:  # pylint: disable=broad-exception-caught
                logger.warning("OCR failed for image %s: %s", filename, exc)

        try:
            text = content.decode('utf-8', errors='ignore')
            lower_text = text[:8]
            if lower_text.startswith("%PDF-") or lower_text.startswith("PK"):
                return ""
            return text
        except Exception as exc:  # pylint: disable=broad-exception-caught
            logger.error("Failed to read file %s: %s", filename, exc)
            return ""


def extract_text_from_file(file_path: Union[str, pathlib.Path]) -> str:
    """Extract text from various file types (txt, pdf, docx, html) or URLs."""
    # Check if it's a URL and normalize single-slash URLs
    file_str = str(file_path)

    if file_str.startswith(('http://', 'https://')):
        content = _download_from_url(file_str)
        if not content:
            return ""
        return extract_text_from_bytes(content, file_str)

    # If not URL, treat as path
    if isinstance(file_path, str):
        file_path = pathlib.Path(file_path)

    if file_path.name in {".DS_Store"}:
        return ""

    # Local file
    try:
        with open(file_path, "rb") as f:
            content = f.read()
        return extract_text_from_bytes(content, file_path.name)
    except Exception as exc:  # pylint: disable=broad-exception-caught
        logger.error("Failed to read file %s: %s", file_path, exc)
        return ""


# Backwards compatibility for legacy imports. Remove once callers migrate.
_extract_text_from_file = extract_text_from_file
